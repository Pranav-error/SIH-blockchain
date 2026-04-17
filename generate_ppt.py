"""
HERB-Vault CMTI DIC 2026 — Comprehensive Presentation Generator v2
Run: python generate_ppt.py
Output: HERB_VAULT_CMTI_DIC2026_FULL.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

GREEN_DARK  = RGBColor(0x1B, 0x5E, 0x20)
GREEN_MID   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0xA5, 0xD6, 0xA7)
GOLD        = RGBColor(0xF9, 0xA8, 0x25)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG    = RGBColor(0xF1, 0xF8, 0xE9)
RED         = RGBColor(0xC6, 0x28, 0x28)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
BLUE        = RGBColor(0x15, 0x65, 0xC0)
BLACK       = RGBColor(0x0D, 0x0D, 0x0D)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── helpers ────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line_col=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_col:
        s.line.color.rgb = line_col; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=14, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def mltxt(slide, items, x, y, w, h, default_size=12,
          default_bold=False, default_color=DARK_TEXT, default_align=PP_ALIGN.LEFT):
    """items = list of str or (text, size, bold, color, align)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            item = (item, default_size, default_bold, default_color, default_align)
        t, s, b, c, a = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = a
        r = p.add_run(); r.text = t
        r.font.size = Pt(s); r.font.bold = b; r.font.color.rgb = c
    return tb

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, H, fill=LIGHT_BG)
    rect(slide, 0, 0, W, Inches(1.0), fill=GREEN_DARK)
    rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill=GOLD)
    txt(slide, title, Inches(0.5), Inches(0.13), Inches(9.5), Inches(0.75),
        size=24, bold=True, color=WHITE)
    txt(slide, "🌿 HERB-VAULT  |  CMTI DIC 2026  |  REVA University",
        W - Inches(4.0), Inches(0.2), Inches(3.8), Inches(0.6),
        size=9, color=GREEN_LIGHT, align=PP_ALIGN.RIGHT)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(0.78), Inches(9), Inches(0.3),
            size=11, color=GREEN_LIGHT, italic=True)

def card_box(slide, x, y, w, h, title, lines,
             head_col=GREEN_MID, head_txt=WHITE, body_col=WHITE,
             head_size=13, body_size=11, txt_color=DARK_TEXT):
    rect(slide, x, y, w, h, fill=body_col, line_col=GREEN_LIGHT, line_w=Pt(1))
    rect(slide, x, y, w, Inches(0.40), fill=head_col)
    txt(slide, title, x+Inches(0.1), y+Inches(0.05), w-Inches(0.15), Inches(0.35),
        size=head_size, bold=True, color=head_txt)
    items = [(l, body_size, False, txt_color, PP_ALIGN.LEFT) for l in lines]
    mltxt(slide, items, x+Inches(0.12), y+Inches(0.47), w-Inches(0.2), h-Inches(0.55))


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, fill=GREEN_DARK)
rect(s, Inches(8.8), 0, Inches(4.53), H, fill=GREEN_MID)
rect(s, 0, H-Inches(0.5), W, Inches(0.5), fill=GOLD)

txt(s,"🌿", Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.1), size=52, color=WHITE)
txt(s,"HERB-VAULT", Inches(0.5), Inches(1.4), Inches(8.1), Inches(1.3),
    size=52, bold=True, color=WHITE)
txt(s,"Blockchain-Based Ayurvedic Herb Traceability System",
    Inches(0.5), Inches(2.65), Inches(8.0), Inches(0.65),
    size=21, color=GREEN_LIGHT)

rect(s, Inches(0.5), Inches(3.45), Inches(3.5), Inches(0.52), fill=GOLD)
txt(s,"📜  PATENT APPLIED — Indian Patent Office",
    Inches(0.62), Inches(3.52), Inches(3.3), Inches(0.4),
    size=12, bold=True, color=GREEN_DARK)

txt(s,"CMTI Design & Innovation Clinic 2026  |  April 16–18, Bengaluru",
    Inches(0.5), Inches(4.1), Inches(8.0), Inches(0.45), size=14, color=GREEN_LIGHT)

rect(s, Inches(0.5), Inches(4.75), Inches(8.0), Inches(2.1), fill=GREEN_MID)
team = [
    ("R Sai Pranav",     "CSE  —  Software Lead"),
    ("Kartik Jarali",    "CSE  —  Software"),
    ("Harshal Y Kumar",  "Mechatronics  —  Hardware Lead"),
    ("Karanam Nayan",    "Mechatronics  —  Hardware"),
]
for i,(name,role) in enumerate(team):
    col_n = i%2; row_n = i//2
    tx = Inches(0.75) + col_n*Inches(3.9)
    ty = Inches(4.9)  + row_n*Inches(0.82)
    txt(s, f"• {name}  |  {role}", tx, ty, Inches(3.7), Inches(0.7),
        size=12, color=WHITE)

txt(s,"INDUSTRY INTEREST", Inches(9.2), Inches(1.4), Inches(3.7), Inches(0.45),
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s,"Himalaya Drug Company\nhas evaluated and expressed\ncommercial interest",
    Inches(9.2), Inches(1.95), Inches(3.7), Inches(1.0),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(9.3), Inches(3.1), Inches(3.5), Inches(0.52), fill=GOLD)
txt(s,"Software Architecture: READY",
    Inches(9.35), Inches(3.17), Inches(3.4), Inches(0.4),
    size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s,"sih-blockchain.vercel.app",
    Inches(9.2), Inches(3.78), Inches(3.7), Inches(0.38),
    size=11, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
txt(s,"Hyperledger Fabric\nBlockchain Network\nDesigned & Tested",
    Inches(9.2), Inches(4.25), Inches(3.7), Inches(1.0),
    size=12, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE CRISIS (Real Loss Data)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"The Real Crisis — Billions Lost, Trust Broken")

# Top stat row
stats = [
    ("₹15,000 Cr+",  "Annual export revenue\nLOST to fraud & rejections",   RED),
    ("47+",          "EU/US border rejections\nof Indian herbs (2023 alone)", RED),
    ("26%",          "Ayurvedic products FAIL\nFSSAI quality sampling",       ORANGE),
    ("40%",          "Herbs in Indian market\nadulterated or mislabeled",     ORANGE),
    ("68%",          "Urban consumers DON'T\ntrust herbal product labels",    BLUE),
    ("3–5×",         "Farm gate vs. shelf price\ngap due to middlemen",       GRAY),
]
for i,(val,label,col) in enumerate(stats):
    bx = Inches(0.3) + i*Inches(2.15)
    rect(s, bx, Inches(1.1), Inches(2.0), Inches(1.65), fill=col)
    txt(s, val, bx+Inches(0.1), Inches(1.18), Inches(1.8), Inches(0.72),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, bx+Inches(0.07), Inches(1.85), Inches(1.88), Inches(0.78),
        size=9, color=WHITE, align=PP_ALIGN.CENTER)

# EU rejection section
rect(s, Inches(0.3), Inches(2.9), Inches(6.2), Inches(4.25), fill=WHITE,
     line_col=RED, line_w=Pt(1.5))
rect(s, Inches(0.3), Inches(2.9), Inches(6.2), Inches(0.42), fill=RED)
txt(s,"🇪🇺  WHY EU & US REJECT INDIAN HERBS",
    Inches(0.45), Inches(2.95), Inches(6.0), Inches(0.35),
    size=13, bold=True, color=WHITE)

eu_points = [
    "Pesticide residues above MRL limits — no field record to prove safe practices",
    "Microbial contamination (Salmonella, E.coli) — no temperature/handling chain record",
    "Incorrect species labeling — Brahmi sold as Gotu Kola, no DNA-verifiable origin",
    "Falsified certificates of origin — paper-based, easy to forge without blockchain proof",
    "Heavy metal contamination — no soil/zone restriction at point of collection",
    "Result: India loses ₹15,000 Cr+ annually in rejected herb exports",
    "EU RASFF portal: 47+ rapid alerts for Indian spices & herbs in 2022–23 alone",
    "Germany, Netherlands, UK are top rejection origin countries for Indian herbs",
]
for i,pt in enumerate(eu_points):
    by = Inches(3.42) + i*Inches(0.46)
    col_c = RED if "Result" in pt or "RASFF" in pt or "Germany" in pt else DARK_TEXT
    b = "Result" in pt
    txt(s, f"{'⚠' if b else '•'}  {pt}", Inches(0.45), by, Inches(5.9), Inches(0.44),
        size=10, bold=b, color=col_c)

# Consumer trust section
rect(s, Inches(6.7), Inches(2.9), Inches(6.3), Inches(4.25), fill=WHITE,
     line_col=BLUE, line_w=Pt(1.5))
rect(s, Inches(6.7), Inches(2.9), Inches(6.3), Inches(0.42), fill=BLUE)
txt(s,"👤  WHY CONSUMERS DON'T TRUST HERBAL MEDICINES",
    Inches(6.85), Inches(2.95), Inches(6.1), Inches(0.35),
    size=13, bold=True, color=WHITE)

trust_points = [
    "\"How do I know this is real Ashwagandha?\" — no way to verify at point of purchase",
    "Multiple fake product scandals: turmeric with lead chromate, curcumin substitution",
    "No visibility into where herbs were grown, by whom, or under what conditions",
    "No lab test results accessible to the buyer — they must trust the brand blindly",
    "Global trend: 72% of consumers pay 15–20% premium for traceable products (Nielsen 2023)",
    "Coffee (Rainforest Alliance QR), wine (blockchain provenance) — premium pricing proven",
    "When a consumer SEES the complete journey, credibility of the product becomes unquestionable",
    "With HERB-Vault QR: farm → collector ID → GPS zone → grade → lab → product",
]
for i,pt in enumerate(trust_points):
    by = Inches(3.42) + i*Inches(0.46)
    b = "When a consumer" in pt or "With HERB-Vault" in pt
    col_c = GREEN_DARK if b else DARK_TEXT
    txt(s, f"{'★' if b else '•'}  {pt}", Inches(6.85), by, Inches(6.0), Inches(0.44),
        size=10, bold=b, color=col_c)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — QR CONSUMER TRUST JOURNEY
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"How QR Traceability Rebuilds Consumer Trust",
       "Scan once — see everything. Farm to pharmacy, immutably on blockchain.")

rect(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.7), fill=GREEN_DARK)
txt(s,"THE QR JOURNEY — What a consumer sees when they scan the product QR code",
    Inches(0.5), Inches(1.18), Inches(12.3), Inches(0.55),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

journey = [
    ("📍 COLLECTION\nEVENT",
     "Ashwagandha\nCollected: 14 Apr 2026\nLocation: Nagaur, Rajasthan\nGPS: 27.2°N 73.7°E\nGPS ZONE: ✓ VERIFIED",
     "Collector: COL-Ramesh\nWeight: 12.4 kg\nMoisture: 10.8%\nGrade: A (Premium)\nBlockchain TX: 3f9a2b...",
     GREEN_DARK),
    ("🏭 PROCESSING\nSTEP",
     "Facility: Herbal Pro\nPune, Maharashtra\nProcess: Sun Drying\nDuration: 72 hours\nDate: 18 Apr 2026",
     "Operator: Suresh K\nInput: 12.4 kg\nOutput: 9.8 kg\nLoss (normal): 21%\nBlockchain TX: 8c1d4e...",
     GREEN_MID),
    ("🧪 QUALITY\nTEST",
     "Lab: NABL Certified\nTest Date: 22 Apr 2026\nMoisture: 9.2% ✓\nPesticide: PASS ✓\nHeavy Metals: PASS ✓",
     "Microbial: PASS ✓\nWithanolides: 2.5%\nEU MRL: COMPLIANT\nUS FDA: COMPLIANT\nBlockchain TX: 2e7f9c...",
     BLUE),
    ("💊 FINAL\nPRODUCT",
     "Himalaya Ashwagandha\nCapsules 500mg\nBatch: HIM-ASH-2026-042\nMfg: 25 Apr 2026\nExp: Apr 2028",
     "Certifications:\n• AYUSH GMP\n• ISO 22000\n• EU Organic\nBlockchain TX: 1a4b5c...",
     GOLD),
]

bw = Inches(3.0); bh = Inches(4.9)
for i,(stage,left,right,col) in enumerate(journey):
    bx = Inches(0.35) + i*(bw + Inches(0.12))
    rect(s, bx, Inches(1.95), bw, bh, fill=col)
    txt(s, stage, bx+Inches(0.1), Inches(2.02), bw-Inches(0.15), Inches(0.75),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, bx+Inches(0.1), Inches(2.85), bw-Inches(0.18), Inches(0.02), fill=GOLD)
    txt(s, left, bx+Inches(0.12), Inches(2.95), bw-Inches(0.22), Inches(1.85),
        size=10, color=WHITE)
    rect(s, bx+Inches(0.1), Inches(4.82), bw-Inches(0.18), Inches(0.02), fill=WHITE)
    txt(s, right, bx+Inches(0.12), Inches(4.9), bw-Inches(0.22), Inches(1.85),
        size=10, color=WHITE)
    if i < 3:
        txt(s,"→", bx+bw+Inches(0.01), Inches(4.2), Inches(0.12), Inches(0.55),
            size=22, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

rect(s, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.38), fill=GOLD)
txt(s,
    "★  Every field above is immutable on Hyperledger Fabric blockchain — no one can alter it after submission."
    "  The consumer isn't trusting a brand. They're reading a cryptographic record.",
    Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.32),
    size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET SIZE & COMPETITION
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Market Opportunity & Competitive Landscape")

# Left — market
rect(s, Inches(0.3), Inches(1.1), Inches(6.0), Inches(6.1), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(1.1), Inches(6.0), Inches(0.42), fill=GREEN_MID)
txt(s,"MARKET OPPORTUNITY", Inches(0.45), Inches(1.15), Inches(5.8), Inches(0.35),
    size=14, bold=True, color=WHITE)

market = [
    ("₹50,000 Cr+", "India Ayurvedic domestic market (2024)",    GOLD),
    ("8.5% CAGR",   "Annual growth rate through 2030",            GREEN_MID),
    ("₹15,000 Cr+", "Export potential by 2030 (AYUSH target)",    GREEN_DARK),
    ("$150 Bn",     "Global herbal supplement market by 2026",     BLUE),
    ("2.5 Lakh+",   "Licensed Ayurvedic manufacturers in India",   GRAY),
    ("5 Cr+",       "Farming families in medicinal herb supply",   GRAY),
    ("₹15,000 Cr",  "Annual losses from fraud & rejections (TAM)", RED),
]
for i,(val,label,c) in enumerate(market):
    by = Inches(1.65) + i*Inches(0.78)
    rect(s, Inches(0.5), by, Inches(1.85), Inches(0.65), fill=c)
    txt(s, val, Inches(0.52), by+Inches(0.1), Inches(1.8), Inches(0.5),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, Inches(2.5), by+Inches(0.17), Inches(3.6), Inches(0.42),
        size=11, color=DARK_TEXT)

txt(s,"Who needs HERB-Vault:",
    Inches(0.5), Inches(7.0), Inches(5.6), Inches(0.3),
    size=11, bold=True, color=GREEN_DARK)
txt(s,"Herb farmers  •  Aggregators  •  Processors  •  Quality labs\nHimalaya, Dabur, Patanjali  •  AYUSH / FSSAI  •  EU importers  •  Consumers",
    Inches(0.5), Inches(7.25), Inches(5.6), Inches(0.5),
    size=9, color=GRAY)

# Right — competition
rect(s, Inches(6.5), Inches(1.1), Inches(6.5), Inches(6.1), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(6.5), Inches(1.1), Inches(6.5), Inches(0.42), fill=GREEN_DARK)
txt(s,"COMPETITIVE LANDSCAPE", Inches(6.65), Inches(1.15), Inches(6.3), Inches(0.35),
    size=14, bold=True, color=WHITE)

headers_c = ["Solution", "GPS\nZone", "Block-\nchain", "Grade\nRecord", "EU\nCompliance", "Consumer\nQR"]
cw = [Inches(1.85), Inches(0.82), Inches(0.82), Inches(0.82), Inches(0.92), Inches(1.05)]
cx = [Inches(6.55)]
for w in cw[:-1]: cx.append(cx[-1]+w)

rect(s, Inches(6.55), Inches(1.62), Inches(6.38), Inches(0.45), fill=GREEN_DARK)
for j,(h,x) in enumerate(zip(headers_c,cx)):
    txt(s, h, x+Inches(0.04), Inches(1.65), cw[j]-Inches(0.05), Inches(0.42),
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Paper Records",              "✗","✗","✗","✗","✗", True),
    ("Generic ERP (SAP/Tally)",    "✗","✗","✓","✗","✗", True),
    ("AgriChain / TraceX",         "✗","Partial","✓","✗","✗", True),
    ("QR-Label solutions",         "✗","✗","✓","✗","✓", True),
    ("IBM Food Trust",             "✗","✓","✓","✗","Partial", True),
    ("HERB-VAULT  ★",              "✓","✓","✓","✓","✓", False),
]
rc = [WHITE,WHITE,WHITE,WHITE,WHITE, RGBColor(0xE8,0xF5,0xE9)]
tm = {"✓":GREEN_MID,"✗":RED,"Partial":GOLD}

for i,(name,g,b,gr,eu,qr,is_comp) in enumerate(rows):
    by = Inches(2.12) + i*Inches(0.78)
    rect(s, Inches(6.55), by, Inches(6.38), Inches(0.74),
         fill=rc[i], line_col=GREEN_LIGHT, line_w=Pt(0.5))
    vals = [name,g,b,gr,eu,qr]
    for j,(v,x) in enumerate(zip(vals,cx)):
        c2 = tm.get(v, GREEN_DARK if not is_comp and j==0 else DARK_TEXT)
        txt(s, v, x+Inches(0.04), by+Inches(0.2), cw[j]-Inches(0.06), Inches(0.45),
            size=10 if j>0 else 9, bold=(not is_comp and j==0) or v in ("✓","✗"),
            color=c2,
            align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

txt(s,"★  HERB-Vault is the ONLY solution with GPS zone enforcement on blockchain AND EU compliance readiness",
    Inches(6.55), Inches(6.95), Inches(6.35), Inches(0.35),
    size=9, bold=True, color=GREEN_DARK)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"HERB-Vault — The Solution")

txt(s,
    "One device + one platform = tamper-proof herb certification from the moment of collection.",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.38),
    size=14, color=GREEN_DARK, italic=True)

pillars = [
    ("🌐 Blockchain\nLedger",
     ["Hyperledger Fabric — permissioned,",
      "enterprise-grade blockchain",
      "Every herb batch gets a permanent,",
      "immutable record across 2 orgs",
      "No single party controls the data"], GREEN_DARK),
    ("📍 GPS Zone\nValidation (Patent)",
     ["Collection GPS validated INSIDE",
      "the smart contract — not server",
      "Haversine formula on-chain",
      "Fake location = blockchain rejects",
      "Patent applied at Indian Patent Office"], GOLD),
    ("📱 Field\nCollector Device",
     ["ESP32-S3 reads weight + moisture",
      "Neo-6M captures live GPS fix",
      "Trained collector selects A/B/C grade",
      "OLED shows result in seconds",
      "No smartphone needed in the field"], GREEN_MID),
    ("📊 Full\nSupply Chain",
     ["Processing facility logs on mobile app",
      "Quality lab submits test results",
      "Manufacturer links batch to product",
      "Consumer scans QR — sees everything",
      "Himalaya audits entire chain"], BLUE),
    ("🇪🇺 Export\nCompliance",
     ["GPS zone + lab test records on chain",
      "EU RASFF-ready audit trail",
      "FSSAI, AYUSH, EU GACP aligned",
      "Pesticide & contaminant records",
      "Eliminates rejection root causes"], GREEN_DARK),
]
bw = Inches(2.45); bh = Inches(4.3)
for i,(title,lines,col) in enumerate(pillars):
    bx = Inches(0.32) + i*(bw+Inches(0.08))
    rect(s, bx, Inches(1.55), bw, bh, fill=col)
    txt(s, title, bx+Inches(0.1), Inches(1.62), bw-Inches(0.15), Inches(0.8),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, bx+Inches(0.1), Inches(2.45), bw-Inches(0.18), Inches(0.02), fill=WHITE)
    for j,line in enumerate(lines):
        txt(s, "• "+line, bx+Inches(0.1), Inches(2.55)+j*Inches(0.6),
            bw-Inches(0.18), Inches(0.58), size=10, color=WHITE)

rect(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(1.15), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
txt(s,"WHY THIS WORKS",
    Inches(0.5), Inches(6.05), Inches(3.5), Inches(0.38),
    size=13, bold=True, color=GREEN_DARK)
why = [
    "• Farmer cannot be cheated — weight and grade are locked to their identity on blockchain",
    "• Middlemen cannot change grades — it's on-chain before they touch the batch",
    "• Pharma company (Himalaya) gets cryptographic proof of source — no paper trust",
    "• Consumer scans QR and sees the complete journey — credibility becomes unquestionable",
    "• Regulators get automated audit trail — no manual inspection of paper records",
]
mltxt(s, [(w, 10, False, DARK_TEXT, PP_ALIGN.LEFT) for w in why],
      Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.95))


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — PIN-TO-PIN WORKFLOW
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Complete System Workflow — Pin to Pin")

steps8 = [
    ("1","COLLECTION\n(Farm/Forest)","Collector places\nherb batch on\ndevice platform","Hardware\nDevice",GREEN_DARK),
    ("2","SENSOR READ\n(Automatic)","Weight via HX711\nMoisture %\nGPS coordinates","IoT Sensors\nAuto-capture",GREEN_MID),
    ("3","GRADE INPUT\n(Collector)","Trained collector\ninspects visually\nPresses A / B / C","3 Buttons\nOn Device",GREEN_MID),
    ("4","GPS VALIDATE\n(Patent)","Haversine check\nruns INSIDE the\nsmart contract","Blockchain\nOn-Chain",GOLD),
    ("5","BLOCKCHAIN\nWRITE","Immutable record:\nweight+moisture\n+grade+GPS+ID","Hyperledger\nFabric",GREEN_DARK),
    ("6","PROCESSING\n& QA","Worker logs steps\nLab submits tests\nvia Mobile App","React Native\nOffline-First",GREEN_MID),
    ("7","FINAL\nPRODUCT","Pharma links batch\nto product, QR\ncode generated","Web Dashboard\nVercel",GREEN_MID),
    ("8","CONSUMER\nVERIFY","Scan QR on phone\nSee full history\nTrust verified","Public QR\nTrace Page",GREEN_DARK),
]
bw = Inches(1.52); bh = Inches(2.6); gap = Inches(0.1)
sx = Inches(0.22); sy = Inches(1.12)
for i,(n,ti,de,tech,c) in enumerate(steps8):
    bx = sx + i*(bw+gap)
    rect(s, bx, sy, bw, bh, fill=c)
    rect(s, bx+Inches(0.55), sy-Inches(0.28), Inches(0.42), Inches(0.42), fill=GOLD)
    txt(s, n, bx+Inches(0.56), sy-Inches(0.27), Inches(0.4), Inches(0.4),
        size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(s, ti, bx+Inches(0.06), sy+Inches(0.08), bw-Inches(0.1), Inches(0.72),
        size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, de, bx+Inches(0.07), sy+Inches(0.85), bw-Inches(0.1), Inches(1.0),
        size=8, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, bx+Inches(0.05), sy+bh-Inches(0.6), bw-Inches(0.08), Inches(0.52), fill=BLACK)
    txt(s, tech, bx+Inches(0.07), sy+bh-Inches(0.58), bw-Inches(0.1), Inches(0.5),
        size=7, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    if i < 7:
        txt(s,"→", bx+bw+Inches(0.005), sy+Inches(0.95), gap+Inches(0.08), Inches(0.45),
            size=16, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

# Role table
rect(s, Inches(0.3), Inches(3.9), Inches(12.7), Inches(3.3), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(3.9), Inches(12.7), Inches(0.4), fill=GREEN_DARK)
txt(s,"STAKEHOLDER ROLE MAP — Who Uses Which Tool",
    Inches(0.5), Inches(3.95), Inches(12.0), Inches(0.32),
    size=13, bold=True, color=WHITE)

roles = [
    ("🌾 Herb Collector",      "Hardware Device\n(ESP32 + sensors + grade buttons)",
     "Submits: weight, moisture, GPS, grade\nNo smartphone required in field"),
    ("🏭 Processing Worker",   "Mobile App\n(Android/iOS, offline-first)",
     "Logs: drying, grinding, storage\nWorks without internet"),
    ("🧪 Quality Lab Tech",    "Mobile App\n(Same app, different role)",
     "Submits: moisture %, pesticide, microbial\nNABL-certified results on chain"),
    ("💊 Pharma (Himalaya)",   "Web Dashboard\n(sih-blockchain.vercel.app)",
     "Views: full audit trail, analytics\nExport compliance reports"),
    ("👤 Consumer",            "QR Code Scan\n(Any phone, no install)",
     "Sees: complete journey farm→shelf\nBuilds product credibility"),
    ("🏛 AYUSH / FSSAI",      "Web Dashboard\n(Regulator access)",
     "Views: compliance reports\nAudit trail for certification"),
]
for i,(who,tool,desc) in enumerate(roles):
    col_n = i%3; row_n = i//2
    rx = Inches(0.45) + col_n*Inches(4.2)
    ry = Inches(4.42) + row_n*Inches(1.3)
    rect(s, rx, ry, Inches(4.0), Inches(1.22), fill=LIGHT_BG,
         line_col=GREEN_LIGHT, line_w=Pt(1))
    txt(s, who, rx+Inches(0.1), ry+Inches(0.06), Inches(3.8), Inches(0.36),
        size=12, bold=True, color=GREEN_DARK)
    txt(s, tool, rx+Inches(0.1), ry+Inches(0.42), Inches(3.8), Inches(0.38),
        size=10, bold=True, color=GREEN_MID)
    txt(s, desc, rx+Inches(0.1), ry+Inches(0.78), Inches(3.8), Inches(0.4),
        size=9, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — HARDWARE DEVICE
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Hardware Device — Field Collector Module",
       "No smartphone. No app. No internet dependency. Just the device.")

rect(s, Inches(0.3), Inches(1.1), Inches(5.8), Inches(6.1), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(1.1), Inches(5.8), Inches(0.42), fill=GREEN_MID)
txt(s,"COMPONENTS", Inches(0.45), Inches(1.15), Inches(5.6), Inches(0.35),
    size=14, bold=True, color=WHITE)

comps = [
    ("ESP32-S3",             "Main controller — WiFi, processing, API calls"),
    ("Neo-6M GPS Module",    "Live GPS coordinates — location proof (UART)"),
    ("HX711 + Load Cell",    "Weighs herb batch accurately"),
    ("Moisture Sensor",      "Reads moisture % of herbs (analog pin)"),
    ("OLED Display 0.96\"",  "Shows grade, batch ID, status to collector"),
    ("3 Grade Buttons A/B/C","Collector presses after visual inspection"),
    ("Green + Red LEDs",     "Instant accepted / rejected visual feedback"),
    ("Power Bank (USB-C)",   "Portable power — works in remote fields"),
]
for i,(comp,desc) in enumerate(comps):
    by = Inches(1.65) + i*Inches(0.61)
    rect(s, Inches(0.45), by, Inches(1.8), Inches(0.52), fill=GREEN_DARK)
    txt(s, comp, Inches(0.5), by+Inches(0.08), Inches(1.7), Inches(0.42),
        size=9, bold=True, color=GOLD)
    txt(s, desc, Inches(2.38), by+Inches(0.1), Inches(3.55), Inches(0.42),
        size=10, color=DARK_TEXT)

rect(s, Inches(6.3), Inches(1.1), Inches(6.7), Inches(6.1), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(6.3), Inches(1.1), Inches(6.7), Inches(0.42), fill=GREEN_DARK)
txt(s,"COLLECTOR STEP-BY-STEP WORKFLOW", Inches(6.45), Inches(1.15),
    Inches(6.5), Inches(0.35), size=14, bold=True, color=WHITE)

csteps = [
    ("1","Place herb batch on device platform"),
    ("2","HX711 auto-reads weight in grams"),
    ("3","Moisture sensor auto-reads %"),
    ("4","Neo-6M captures live GPS coordinates"),
    ("5","Collector visually inspects herb quality"),
    ("6","Collector presses  [A]  [B]  or  [C]  grade button"),
    ("7","ESP32 sends all data to POST /api/intake"),
    ("8","Backend validates GPS zone (patent feature)"),
    ("9","If valid: blockchain record written + batch ID"),
    ("10","OLED: ACCEPTED  |  Grade A  |  Batch ID"),
    ("11","Green LED on — collector hands batch forward"),
    ("12","If rejected: Red LED + OLED shows reason"),
]
for i,(n,step) in enumerate(csteps):
    by = Inches(1.65) + i*Inches(0.50)
    rect(s, Inches(6.45), by, Inches(0.38), Inches(0.42), fill=GOLD)
    txt(s, n, Inches(6.46), by+Inches(0.06), Inches(0.36), Inches(0.36),
        size=9, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(s, step, Inches(6.95), by+Inches(0.07), Inches(5.9), Inches(0.42),
        size=11, color=DARK_TEXT)

# OLED mockup
rect(s, Inches(9.85), Inches(5.6), Inches(2.85), Inches(1.55), fill=BLACK)
rect(s, Inches(9.96), Inches(5.72), Inches(2.62), Inches(1.3), fill=RGBColor(0x04,0x04,0x28))
txt(s,"✓ ACCEPTED", Inches(9.98), Inches(5.77), Inches(2.58), Inches(0.38),
    size=13, bold=True, color=RGBColor(0x00,0xFF,0x88), align=PP_ALIGN.CENTER)
txt(s,"Grade: A  (Premium)", Inches(9.98), Inches(6.1), Inches(2.58), Inches(0.3),
    size=10, color=GOLD, align=PP_ALIGN.CENTER)
txt(s,"ASH-A-20260416-3F9A2B", Inches(9.98), Inches(6.38), Inches(2.58), Inches(0.28),
    size=8, color=WHITE, align=PP_ALIGN.CENTER)
txt(s,"OLED Display Output", Inches(9.9), Inches(7.04), Inches(2.8), Inches(0.28),
    size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — PATENT INNOVATION
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Patent Applied — The Core Innovation")

rect(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.85), fill=GOLD)
txt(s,
    "📜  PATENT APPLICATION FILED — Indian Patent Office"
    "  |  Innovation: GPS Zone Validation Executed Inside Blockchain Smart Contract",
    Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.58),
    size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

txt(s,"Important: The patent protects the METHOD and INNOVATION — not the prototype.",
    Inches(0.5), Inches(2.05), Inches(12.3), Inches(0.38),
    size=12, bold=True, color=RED, italic=True)
txt(s,"The idea of performing Haversine GPS validation as a smart contract transaction is patent-protected regardless of the hardware used.",
    Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.38),
    size=11, color=DARK_TEXT)

rect(s, Inches(0.3), Inches(2.88), Inches(5.9), Inches(4.3), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(2.88), Inches(5.9), Inches(0.4), fill=GREEN_DARK)
txt(s,"WHAT THE PATENT PROTECTS", Inches(0.45), Inches(2.92),
    Inches(5.7), Inches(0.32), size=13, bold=True, color=WHITE)
pat = [
    "Collection GPS validated INSIDE the smart contract, not the application server",
    "Haversine formula calculates great-circle distance within chaincode execution",
    "If collector is outside approved zone → blockchain REJECTS at consensus level",
    "No backend bypass possible — all network nodes enforce the rule simultaneously",
    "Each herb species has species-specific approved zones encoded in the contract",
    "Collector identity, GPS proof, and grade are cryptographically linked in one TX",
    "First known on-chain geospatial herb zone enforcement in permissioned blockchain",
]
for i,p in enumerate(pat):
    by = Inches(3.38) + i*Inches(0.55)
    rect(s, Inches(0.45), by+Inches(0.1), Inches(0.26), Inches(0.26), fill=GOLD)
    txt(s,"✓", Inches(0.46), by+Inches(0.08), Inches(0.24), Inches(0.28),
        size=10, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(s, p, Inches(0.8), by+Inches(0.04), Inches(5.25), Inches(0.5),
        size=10, color=DARK_TEXT)

rect(s, Inches(6.5), Inches(2.88), Inches(6.5), Inches(4.3), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(6.5), Inches(2.88), Inches(6.5), Inches(0.4), fill=GREEN_MID)
txt(s,"BLOCKCHAIN ARCHITECTURE", Inches(6.65), Inches(2.92),
    Inches(6.3), Inches(0.32), size=13, bold=True, color=WHITE)
arch = [
    ("Platform",    "Hyperledger Fabric 2.4.x — enterprise permissioned blockchain"),
    ("Channel",     "herblock — dedicated Ayurvedic herb supply chain"),
    ("Orgs",        "Org1MSP (Govt/Regulator) + Org2MSP (Industry) — both must endorse"),
    ("Consensus",   "Raft (Crash Fault Tolerant) — no single point of failure"),
    ("Chaincode",   "Node.js — GPS Haversine validation executes here (patent)"),
    ("World State", "CouchDB — queryable ledger state + MongoDB app cache"),
    ("Security",    "TLS 1.3 on all peer communication"),
    ("Hashing",     "SHA-256 fingerprint per record + Merkle root per batch"),
    ("Immutability","No record can be altered after commitment — by anyone"),
]
for i,(k,v) in enumerate(arch):
    by = Inches(3.38) + i*Inches(0.48)
    txt(s, f"{k}:", Inches(6.65), by, Inches(1.5), Inches(0.44),
        size=10, bold=True, color=GREEN_DARK)
    txt(s, v, Inches(8.2), by, Inches(4.6), Inches(0.44), size=10, color=DARK_TEXT)

rect(s, Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.22), fill=GREEN_DARK)
txt(s,"No competitor can replicate this — the patent protection covers the method, not just the implementation.",
    Inches(0.5), Inches(7.22), Inches(12.3), Inches(0.2),
    size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — SOFTWARE ARCHITECTURE READY
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Software Architecture — Designed & Ready",
       "The full platform is architected, built, and tested. CMTI is where hardware meets it.")

rect(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(0.65), fill=GOLD)
txt(s,
    "We did not bring a ready-made product. We built the software architecture so the hardware team "
    "has a working platform to integrate with — from day one of the build.",
    Inches(0.5), Inches(1.16), Inches(12.3), Inches(0.54),
    size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

layers = [
    ("CLIENT LAYER",
     ["React 19 Web Dashboard — sih-blockchain.vercel.app",
      "React Native Mobile App — Android/iOS, offline-first SQLite",
      "Public QR Trace Page — no install, any phone",
      "Hardware Device UI — OLED + LEDs on ESP32"],
     BLUE, Inches(0.3), Inches(1.88)),
    ("API LAYER",
     ["FastAPI (Python) — REST endpoints for all supply chain events",
      "POST /api/intake — hardware device endpoint (built for CMTI)",
      "JWT + Google OAuth2 authentication",
      "GPS zone validation, QR generation, analytics dashboard"],
     GREEN_MID, Inches(0.3), Inches(3.28)),
    ("BLOCKCHAIN LAYER",
     ["Hyperledger Fabric 2.4.x — herblock channel",
      "Node.js Chaincode — GPS Haversine validation (patent applied)",
      "2 Orgs (Govt + Industry) — Raft consensus endorsement",
      "Immutable records: CollectionEvent, ProcessingStep, QualityTest, Product"],
     GREEN_DARK, Inches(0.3), Inches(4.68)),
    ("DATA LAYER",
     ["MongoDB — application data (fast queries, analytics)",
      "CouchDB — Fabric world state (blockchain ledger state)",
      "Expo SQLite — mobile offline storage",
      "SHA-256 hash chain + Merkle root per batch"],
     GRAY, Inches(0.3), Inches(6.08)),
]
for title,lines,col,lx,ly in layers:
    rect(s, lx, ly, Inches(12.7), Inches(1.25), fill=col)
    txt(s, title, lx+Inches(0.15), ly+Inches(0.08), Inches(2.5), Inches(0.4),
        size=11, bold=True, color=GOLD)
    for j,line in enumerate(lines):
        col_n = j%2; row_n = j//2
        tx = lx + Inches(2.8) + col_n*Inches(4.9)
        ty = ly + Inches(0.12) + row_n*Inches(0.52)
        txt(s, "• "+line, tx, ty, Inches(4.8), Inches(0.48), size=10, color=WHITE)

# Arrows between layers
for ay in [Inches(3.18), Inches(4.58), Inches(5.98)]:
    txt(s,"↕", Inches(6.4), ay, Inches(0.5), Inches(0.22),
        size=14, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

txt(s,"← Hardware device connects here via POST /api/intake with X-Device-Key auth",
    Inches(5.5), Inches(3.98), Inches(7.5), Inches(0.28),
    size=9, bold=True, color=GOLD, italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — WHAT WE BUILD AT CMTI
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"What We Are Building at CMTI — April 16–18")

rect(s, Inches(0.3), Inches(1.1), Inches(6.0), Inches(6.1), fill=WHITE,
     line_col=GRAY, line_w=Pt(1))
rect(s, Inches(0.3), Inches(1.1), Inches(6.0), Inches(0.42), fill=GRAY)
txt(s,"SOFTWARE ARCHITECTURE  (came in with this)",
    Inches(0.45), Inches(1.15), Inches(5.8), Inches(0.35),
    size=12, bold=True, color=WHITE)
txt(s,"This is our platform — like a startup's app before it launches.",
    Inches(0.45), Inches(1.6), Inches(5.7), Inches(0.35),
    size=9, color=GRAY, italic=True)
prebuilt = [
    "Hyperledger Fabric blockchain network — 2 orgs, Raft",
    "Node.js chaincode with on-chain GPS validation (patent applied)",
    "FastAPI backend — all supply chain REST endpoints",
    "POST /api/intake endpoint — ready for hardware connection",
    "React 19 web dashboard — live on Vercel",
    "React Native mobile app — offline-first, SQLite sync",
    "JWT + Google OAuth2 authentication",
    "QR code generation + consumer trace page",
    "Analytics dashboard with real-time statistics",
]
for i,item in enumerate(prebuilt):
    txt(s, "✓  "+item, Inches(0.5), Inches(2.05)+i*Inches(0.55),
        Inches(5.6), Inches(0.52), size=10, color=GRAY)

rect(s, Inches(6.5), Inches(1.1), Inches(6.5), Inches(6.1), fill=WHITE,
     line_col=GOLD, line_w=Pt(2))
rect(s, Inches(6.5), Inches(1.1), Inches(6.5), Inches(0.42), fill=GOLD)
txt(s,"🔧  BUILDING HERE  (hardware — the last mile)",
    Inches(6.65), Inches(1.15), Inches(6.3), Inches(0.35),
    size=12, bold=True, color=GREEN_DARK)
txt(s,"This is what completes the product. Built from scratch at CMTI.",
    Inches(6.65), Inches(1.6), Inches(6.2), Inches(0.35),
    size=9, color=GREEN_DARK, italic=True, bold=True)
building = [
    "ESP32-S3 firmware — written from scratch here",
    "HX711 load cell wiring + weight calibration",
    "Moisture sensor wiring + analog reading",
    "Neo-6M GPS module integration + live satellite fix",
    "3-button grade selection logic (A / B / C)",
    "OLED display output — grade + batch ID + status",
    "Green / Red LED feedback circuit",
    "WiFi connection + HTTP POST to /api/intake",
    "End-to-end: device → API → blockchain → dashboard",
]
for i,item in enumerate(building):
    rect(s, Inches(6.65), Inches(2.05)+i*Inches(0.55), Inches(0.3), Inches(0.3), fill=GOLD)
    txt(s,"★", Inches(6.66), Inches(2.04)+i*Inches(0.55), Inches(0.28), Inches(0.32),
        size=10, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(s, item, Inches(7.08), Inches(2.05)+i*Inches(0.55),
        Inches(5.75), Inches(0.52), size=11, bold=True, color=GREEN_DARK)

rect(s, Inches(6.5), Inches(7.0), Inches(6.5), Inches(0.42), fill=GREEN_DARK)
txt(s,"The hardware + firmware integration IS the build. It happens here, live.",
    Inches(6.65), Inches(7.04), Inches(6.3), Inches(0.35),
    size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL & VALIDATION
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Business Model & Industry Validation")

rect(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(1.05), fill=GOLD)
txt(s,"🏆  HIMALAYA DRUG COMPANY",
    Inches(0.5), Inches(1.16), Inches(3.8), Inches(0.45),
    size=17, bold=True, color=GREEN_DARK)
txt(s,"Evaluated HERB-Vault and expressed commercial interest.\n"
      "Request: \"Show us working hardware.\"  —  That is exactly what we are building at CMTI.",
    Inches(4.4), Inches(1.18), Inches(8.4), Inches(0.88),
    size=12, color=GREEN_DARK)

rect(s, Inches(0.3), Inches(2.28), Inches(12.7), Inches(0.38), fill=GREEN_DARK)
txt(s,"REVENUE MODEL", Inches(0.5), Inches(2.32), Inches(12.0), Inches(0.3),
    size=13, bold=True, color=WHITE)
revenue = [
    ("Device Sales",    "₹8,000–15,000\nper unit",     "Hardware device sold to\ncooperatives / pharma cos"),
    ("SaaS Dashboard",  "₹2,000–5,000\n/month/org",    "Web + API access for\nmanufacturers"),
    ("Per-TX Fee",      "₹2–5\nper batch record",       "Blockchain write fee\ncharged to processor"),
    ("Compliance Data", "₹50K–2L/year",                 "Harvest & audit reports\nfor AYUSH/exporters"),
]
for i,(title,price,desc) in enumerate(revenue):
    bx = Inches(0.4)+i*Inches(3.1)
    rect(s, bx, Inches(2.78), Inches(2.9), Inches(1.75), fill=WHITE,
         line_col=GREEN_LIGHT, line_w=Pt(1))
    rect(s, bx, Inches(2.78), Inches(2.9), Inches(0.38), fill=GREEN_MID)
    txt(s, title, bx+Inches(0.1), Inches(2.81), Inches(2.7), Inches(0.32),
        size=12, bold=True, color=WHITE)
    txt(s, price, bx+Inches(0.1), Inches(3.22), Inches(2.7), Inches(0.62),
        size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, desc, bx+Inches(0.1), Inches(3.85), Inches(2.7), Inches(0.58),
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.3), Inches(4.65), Inches(6.0), Inches(2.55), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(4.65), Inches(6.0), Inches(0.38), fill=GREEN_MID)
txt(s,"GO-TO-MARKET", Inches(0.45), Inches(4.69), Inches(5.8), Inches(0.32),
    size=13, bold=True, color=WHITE)
gtm = [
    "Phase 1 (Now):   Pilot — 10 collectors, 1 herb zone, with Himalaya",
    "Phase 2 (6 mo):  5 herb species, 3 processors, mobile app rollout",
    "Phase 3 (1 yr):  License platform to Dabur, Patanjali, export labs",
    "Phase 4 (2 yr):  EU/US compliance module — capture export market",
]
for i,g in enumerate(gtm):
    txt(s, "→  "+g, Inches(0.5), Inches(5.14)+i*Inches(0.48),
        Inches(5.7), Inches(0.44), size=10, color=DARK_TEXT)

rect(s, Inches(6.5), Inches(4.65), Inches(6.5), Inches(2.55), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(6.5), Inches(4.65), Inches(6.5), Inches(0.38), fill=GREEN_DARK)
txt(s,"WHY HERB-VAULT WINS", Inches(6.65), Inches(4.69), Inches(6.3), Inches(0.32),
    size=13, bold=True, color=WHITE)
why = [
    "Only blockchain + hardware integration targeting Ayurveda sector",
    "Patent applied — GPS on-chain method cannot be copied legally",
    "Solves EU rejection root cause — verifiable audit trail from farm",
    "Consumer QR trust — credibility built by transparency, not branding",
    "Himalaya validation — real commercial interest, not hypothetical",
    "₹50,000 Cr market + ₹15,000 Cr annual loss = proven, measurable need",
]
for i,w in enumerate(why):
    txt(s, "★  "+w, Inches(6.7), Inches(5.14)+i*Inches(0.38),
        Inches(6.1), Inches(0.36), size=10, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — LIVE DEMO PLAN
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s,"Live Demo — What Judges Will See")

demo = [
    ("STEP 1","Place herb\nsample on\ndevice platform",
     "Collector places a\nbag of herbs on\nthe load cell",GREEN_DARK),
    ("STEP 2","Sensors\nauto-capture\nreadings",
     "Weight: HX711\nMoisture: sensor\nGPS: Neo-6M fix",GREEN_MID),
    ("STEP 3","Collector\nselects grade\nA / B / C",
     "Visual inspection\nPress grade button\non the device",GREEN_MID),
    ("STEP 4","ESP32 POSTs\nto /api/intake\n+ Device Key",
     "JSON: weight +\nmoisture + GPS\n+ grade sent",GOLD),
    ("STEP 5","GPS validated\nON blockchain\n(patent live)",
     "Haversine runs\ninside chaincode\nNot our server",GREEN_DARK),
    ("STEP 6","OLED + LED\nfeedback to\ncollector",
     "ACCEPTED\nGrade: A\nBatch ID shown",GREEN_MID),
    ("STEP 7","Dashboard\nupdates in\nreal time",
     "New blockchain\nrecord appears\nOpen the laptop",GREEN_MID),
    ("STEP 8","Consumer\nscans QR\non phone",
     "Full journey:\nfarm → lab → shelf\nTrust verified",GREEN_DARK),
]
bw = Inches(1.52); bh = Inches(2.65); gap2 = Inches(0.1)
sx2 = Inches(0.22)
for i,(step,title,detail,col) in enumerate(demo):
    bx = sx2 + i*(bw+gap2)
    rect(s, bx, Inches(1.12), bw, bh, fill=col)
    rect(s, bx, Inches(1.12), bw, Inches(0.38), fill=BLACK)
    txt(s, step, bx+Inches(0.06), Inches(1.16), bw-Inches(0.1), Inches(0.32),
        size=9, bold=True, color=GOLD)
    txt(s, title, bx+Inches(0.06), Inches(1.55), bw-Inches(0.1), Inches(0.78),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, detail, bx+Inches(0.07), Inches(2.38), bw-Inches(0.1), Inches(1.3),
        size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 7:
        txt(s,"→", bx+bw+Inches(0.005), Inches(1.9), gap2+Inches(0.08), Inches(0.45),
            size=15, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

rect(s, Inches(0.3), Inches(3.95), Inches(12.7), Inches(3.25), fill=WHITE,
     line_col=GREEN_LIGHT, line_w=Pt(1))
rect(s, Inches(0.3), Inches(3.95), Inches(12.7), Inches(0.38), fill=GREEN_DARK)
txt(s,"JUDGE Q&A — EXPECTED QUESTIONS & ANSWERS",
    Inches(0.5), Inches(3.99), Inches(12.0), Inches(0.32),
    size=13, bold=True, color=WHITE)
qa = [
    ("Q: Isn't this already built?",
     "A: The software architecture is ready — like a startup's platform before launch. "
     "The hardware device and its full firmware integration is being built live here at CMTI."),
    ("Q: How is this different from a QR label?",
     "A: A QR label is a sticker. HERB-Vault is a cryptographic proof. "
     "The GPS, weight, grade and collector identity are on a tamper-proof blockchain — no one can change it."),
    ("Q: Can someone fake the GPS?",
     "A: The GPS validation runs inside the Hyperledger Fabric smart contract — not our server. "
     "All network nodes enforce it simultaneously. This is the patent. It cannot be bypassed."),
    ("Q: Who will pay for this?",
     "A: Himalaya Drug Company has already shown interest. They pay for supply chain certification "
     "today via paper — we replace that with a ₹2–5 per batch blockchain record."),
]
for i,(q,a) in enumerate(qa):
    col_n = i%2; row_n = i//2
    qx = Inches(0.45)+col_n*Inches(6.3)
    qy = Inches(4.45)+row_n*Inches(1.25)
    rect(s, qx, qy, Inches(6.1), Inches(1.18), fill=LIGHT_BG,
         line_col=GREEN_LIGHT, line_w=Pt(1))
    txt(s, q, qx+Inches(0.1), qy+Inches(0.06), Inches(5.9), Inches(0.35),
        size=11, bold=True, color=GREEN_DARK)
    txt(s, a, qx+Inches(0.1), qy+Inches(0.42), Inches(5.9), Inches(0.7),
        size=10, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, fill=GREEN_DARK)
rect(s, Inches(8.9), 0, Inches(4.43), H, fill=GREEN_MID)
rect(s, 0, H-Inches(0.5), W, Inches(0.5), fill=GOLD)

txt(s,"🌿  HERB-VAULT", Inches(0.5), Inches(0.45), Inches(8.3), Inches(1.1),
    size=46, bold=True, color=WHITE)
txt(s,"From Farm to Pharmacy — Trustless, Transparent, Traceable",
    Inches(0.5), Inches(1.5), Inches(8.3), Inches(0.6),
    size=19, color=GREEN_LIGHT)

pillars3 = [
    ("PATENT\nAPPLIED","GPS On-Chain Validation\nIndian Patent Office",GOLD),
    ("HIMALAYA\nINTEREST","Commercial interest expressed\nAwaiting hardware demo",GREEN_MID),
    ("₹15,000 Cr\nANNUAL LOSS","To fraud & rejections\nWe solve the root cause",RED),
]
for i,(ti,su,c) in enumerate(pillars3):
    bx = Inches(0.5)+i*Inches(2.65)
    rect(s, bx, Inches(2.3), Inches(2.4), Inches(2.0), fill=c)
    txt(s, ti, bx+Inches(0.12), Inches(2.42), Inches(2.2), Inches(0.75),
        size=18, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    txt(s, su, bx+Inches(0.12), Inches(3.15), Inches(2.2), Inches(0.65),
        size=12, color=GREEN_DARK, align=PP_ALIGN.CENTER)

txt(s,
    "The Ayurvedic herb supply chain loses ₹15,000 Cr annually\n"
    "because there is no verifiable record at point of collection.\n"
    "HERB-Vault is the only solution that enforces authenticity\n"
    "with patented on-chain GPS validation — backed by hardware.",
    Inches(0.5), Inches(4.55), Inches(8.0), Inches(1.7),
    size=14, color=WHITE)

txt(s,"Software Architecture: READY  |  Patent: APPLIED  |  Industry: INTERESTED",
    Inches(0.5), Inches(6.25), Inches(8.0), Inches(0.45),
    size=12, bold=True, color=GOLD)
txt(s,"github.com/Pranav-error/SIH-blockchain  |  REVA University  |  CMTI DIC 2026",
    Inches(0.5), Inches(6.75), Inches(8.0), Inches(0.38),
    size=11, color=GREEN_LIGHT)

txt(s,"WHAT WE'RE\nBUILDING HERE", Inches(9.2), Inches(1.2), Inches(3.7), Inches(0.85),
    size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
building2 = [
    "ESP32-S3 firmware",
    "HX711 weight calibration",
    "Moisture sensor reading",
    "Neo-6M GPS integration",
    "A/B/C grade button logic",
    "OLED display output",
    "Full API connection",
    "End-to-end demo loop",
]
for i,b in enumerate(building2):
    txt(s, "★  "+b, Inches(9.2), Inches(2.2)+i*Inches(0.52), Inches(3.7), Inches(0.48),
        size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — SYSTEM ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, fill=RGBColor(0xF4, 0xF9, 0xF4))
rect(s, 0, 0, W, Inches(0.92), fill=GREEN_DARK)
rect(s, 0, H - Inches(0.3), W, Inches(0.3), fill=GOLD)
txt(s, "System Architecture — Complete Data Flow",
    Inches(0.4), Inches(0.1), Inches(10.0), Inches(0.72),
    size=22, bold=True, color=WHITE)
txt(s, "🌿 HERB-VAULT  |  CMTI DIC 2026  |  REVA University",
    W - Inches(4.0), Inches(0.18), Inches(3.8), Inches(0.55),
    size=9, color=GREEN_LIGHT, align=PP_ALIGN.RIGHT)

# ── Layer labels (left side vertical) ────────────────────────
def layer_label(slide, label, y, color):
    rect(slide, Inches(0.0), y, Inches(0.22), Inches(0.52), fill=color)

# ── Helper: draw a thick arrow (right-pointing) ──────────────
def arrow_right(slide, x, y):
    """Draw a → using a filled chevron-ish shape via text"""
    txt(slide, "▶", x, y, Inches(0.35), Inches(0.45),
        size=20, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

def arrow_down(slide, x, y):
    txt(slide, "▼", x, y, Inches(0.45), Inches(0.35),
        size=18, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# ROW 1  —  MAIN DATA INGESTION FLOW  (y ≈ 1.0 – 4.55)
# Three main boxes: Field → Backend → Blockchain
# ═══════════════════════════════════════════════════════

B1X, B1Y, B1W, B1H = Inches(0.25), Inches(1.0), Inches(3.65), Inches(3.5)
B2X, B2Y, B2W, B2H = Inches(4.55), Inches(1.0), Inches(4.05), Inches(3.5)
B3X, B3Y, B3W, B3H = Inches(9.25), Inches(1.0), Inches(3.8),  Inches(3.5)

# ── Box 1: Field Device ──────────────────────────────────────
rect(s, B1X, B1Y, B1W, B1H, fill=WHITE, line_col=GREEN_MID, line_w=Pt(2))
rect(s, B1X, B1Y, B1W, Inches(0.48), fill=GREEN_DARK)
txt(s, "① FIELD DEVICE  (ESP32-S3)",
    B1X+Inches(0.1), B1Y+Inches(0.07), B1W-Inches(0.15), Inches(0.38),
    size=12, bold=True, color=WHITE)

esp_items = [
    ("⚖️", "Weight sensor (HX711 + load cell)"),
    ("💧", "Moisture % (capacitive sensor)"),
    ("📍", "GPS coordinates (Neo-6M module)"),
    ("🅰", "Collector presses Grade  A / B / C"),
    ("🔑", "Device authenticates with API Key"),
    ("📟", "OLED: ACCEPTED / REJECTED + Batch ID"),
    ("🟢🔴", "Green LED (accept)  |  Red LED (reject)"),
]
for i, (icon, label) in enumerate(esp_items):
    ty = B1Y + Inches(0.58) + i * Inches(0.41)
    txt(s, icon, B1X+Inches(0.1), ty, Inches(0.32), Inches(0.38), size=11)
    txt(s, label, B1X+Inches(0.42), ty, B1W-Inches(0.52), Inches(0.38),
        size=10, color=DARK_TEXT)

# API call label on arrow
txt(s, "POST /api/intake\n+ X-Device-Key",
    Inches(3.92), Inches(2.25), Inches(0.62), Inches(0.75),
    size=7.5, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
arrow_right(s, Inches(4.13), Inches(2.48))

# ── Box 2: Backend ───────────────────────────────────────────
rect(s, B2X, B2Y, B2W, B2H, fill=WHITE, line_col=BLUE, line_w=Pt(2))
rect(s, B2X, B2Y, B2W, Inches(0.48), fill=BLUE)
txt(s, "② BACKEND  (FastAPI + MongoDB)",
    B2X+Inches(0.1), B2Y+Inches(0.07), B2W-Inches(0.15), Inches(0.38),
    size=12, bold=True, color=WHITE)

be_items = [
    ("🔑", "Validate Device API Key  →  db.devices"),
    ("📍", "GPS Geo-Fence (Haversine)  ★ PATENT ★"),
    ("📅", "Harvest Season (NMPB calendar)"),
    ("✅", "Compute Grade from collector input"),
    ("🆔", "Generate Batch ID  ASH-A-YYYYMMDD-XXXX"),
    ("💾", "Store event in MongoDB  (collection_events)"),
    ("🔗", "Invoke Fabric chaincode  (recordCollection)"),
    ("📱", "Generate QR code  (base64 PNG)"),
]
for i, (icon, label) in enumerate(be_items):
    ty = B2Y + Inches(0.58) + i * Inches(0.365)
    txt(s, icon, B2X+Inches(0.1), ty, Inches(0.32), Inches(0.35), size=10)
    txt(s, label, B2X+Inches(0.42), ty, B2W-Inches(0.52), Inches(0.35),
        size=9.5, color=DARK_TEXT,
        bold=(label.find("PATENT") != -1))

# Fabric invoke label on arrow
txt(s, "peer chaincode\ninvoke (CLI)",
    Inches(8.68), Inches(2.3), Inches(0.55), Inches(0.6),
    size=7.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
arrow_right(s, Inches(8.9), Inches(2.48))

# ── Box 3: Blockchain ────────────────────────────────────────
rect(s, B3X, B3Y, B3W, B3H, fill=WHITE, line_col=GOLD, line_w=Pt(2))
rect(s, B3X, B3Y, B3W, Inches(0.48), fill=RGBColor(0x7B, 0x35, 0x00))
txt(s, "③ HYPERLEDGER FABRIC",
    B3X+Inches(0.1), B3Y+Inches(0.07), B3W-Inches(0.15), Inches(0.38),
    size=12, bold=True, color=GOLD)

bc_items = [
    ("⛓️", "Channel: herblock"),
    ("🏛️", "Org1MSP  —  Government node"),
    ("🏭", "Org2MSP  —  Industry node"),
    ("✍️", "Both orgs must endorse TX (consensus)"),
    ("📜", "Chaincode: recordCollection (Node.js)"),
    ("📍", "GPS Haversine runs INSIDE chaincode"),
    ("🔒", "SHA-256 hash  →  immutable record"),
    ("🌍", "Raft ordering  |  CouchDB world state"),
]
for i, (icon, label) in enumerate(bc_items):
    ty = B3Y + Inches(0.58) + i * Inches(0.365)
    txt(s, icon, B3X+Inches(0.1), ty, Inches(0.32), Inches(0.35), size=10)
    txt(s, label, B3X+Inches(0.42), ty, B3W-Inches(0.52), Inches(0.35),
        size=9.5, color=DARK_TEXT,
        bold=(label.find("INSIDE") != -1 or label.find("Haversine") != -1))

# ═══════════════════════════════════════════════════════
# DOWN ARROW — Blockchain → Consumer verification row
# ═══════════════════════════════════════════════════════
arrow_down(s, Inches(11.0), Inches(4.55))
txt(s, "batch_id\n+ QR code", Inches(10.95), Inches(4.42),
    Inches(0.85), Inches(0.52), size=7, color=GREEN_DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# ROW 2  —  CONSUMER VERIFICATION + RESPONSE PATHS
# ═══════════════════════════════════════════════════════
ROW2Y = Inches(5.02)
ROW2H = Inches(2.1)

# Box A: Web/Mobile App
rect(s, Inches(0.25), ROW2Y, Inches(3.65), ROW2H,
     fill=WHITE, line_col=GREEN_MID, line_w=Pt(2))
rect(s, Inches(0.25), ROW2Y, Inches(3.65), Inches(0.4), fill=GREEN_MID)
txt(s, "④ WEB / MOBILE APP",
    Inches(0.35), ROW2Y+Inches(0.06), Inches(3.45), Inches(0.32),
    size=11, bold=True, color=WHITE)
app_items = [
    "React 19 web  +  React Native mobile",
    "Consumer scans QR code with phone",
    "GET /api/trace/{batch_id}  →  backend",
    "Full journey: Farm → Lab → Shelf",
    "Displays GPS map + grade + lab results",
]
for i, line in enumerate(app_items):
    txt(s, "→  " + line, Inches(0.35), ROW2Y+Inches(0.50)+i*Inches(0.32),
        Inches(3.5), Inches(0.3), size=9.5, color=DARK_TEXT)

# Arrow: App → Backend (api/trace query)
txt(s, "GET /api/trace\n/{batch_id}", Inches(3.92), ROW2Y+Inches(0.55),
    Inches(0.62), Inches(0.55), size=7.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "◀", Inches(4.14), ROW2Y+Inches(0.82), Inches(0.35), Inches(0.38),
    size=20, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

# Box B: MongoDB (data store)
rect(s, Inches(4.55), ROW2Y, Inches(4.05), ROW2H,
     fill=WHITE, line_col=BLUE, line_w=Pt(2))
rect(s, Inches(4.55), ROW2Y, Inches(4.05), Inches(0.4), fill=BLUE)
txt(s, "⑤ DATA STORE  (MongoDB)",
    Inches(4.65), ROW2Y+Inches(0.06), Inches(3.85), Inches(0.32),
    size=11, bold=True, color=WHITE)
mongo_cols = [
    ("collection_events", "Raw intake: weight, GPS, grade, hash"),
    ("processing_steps",  "Processing stages from mobile app"),
    ("quality_tests",     "Lab test results (QA mobile screen)"),
    ("products",          "Final product + QR code link"),
    ("blockchain_txns",   "Fabric TX id + block number log"),
]
for i, (col, desc) in enumerate(mongo_cols):
    ty = ROW2Y + Inches(0.50) + i * Inches(0.32)
    txt(s, col, Inches(4.65), ty, Inches(1.52), Inches(0.3),
        size=8.5, bold=True, color=BLUE)
    txt(s, desc, Inches(6.2), ty, Inches(2.3), Inches(0.3),
        size=8.5, color=DARK_TEXT)

# Arrow up: Backend down to Mongo
arrow_down(s, Inches(6.35), Inches(4.55))

# Box C: Consumer QR Outcome
rect(s, Inches(9.25), ROW2Y, Inches(3.8), ROW2H,
     fill=WHITE, line_col=GOLD, line_w=Pt(2))
rect(s, Inches(9.25), ROW2Y, Inches(3.8), Inches(0.4),
     fill=RGBColor(0x7B, 0x35, 0x00))
txt(s, "⑥ CONSUMER TRUST",
    Inches(9.35), ROW2Y+Inches(0.06), Inches(3.6), Inches(0.32),
    size=11, bold=True, color=GOLD)
consumer_items = [
    "Batch ID  →  collection GPS verified",
    "Collector ID + Grade (A/B/C) visible",
    "Processing plant + date stamp",
    "Lab quality test results on-chain",
    "Zero trust required — blockchain proof",
]
for i, line in enumerate(consumer_items):
    txt(s, "✓  " + line, Inches(9.35), ROW2Y+Inches(0.50)+i*Inches(0.32),
        Inches(3.6), Inches(0.3), size=9.5, color=DARK_TEXT)

# ── Vertical connecting line Blockchain → Consumer box ───────
# already covered by arrow_down above

# ── Left-pointing arrow: Consumer box → Web App ──────────────
txt(s, "◀", Inches(8.82), ROW2Y+Inches(0.82), Inches(0.38), Inches(0.38),
    size=20, bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)
txt(s, "QR scan\nresult", Inches(8.68), ROW2Y+Inches(0.55),
    Inches(0.55), Inches(0.5), size=7.5, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ── Patent badge overlay on GPS validation step ───────────────
rect(s, Inches(4.55), Inches(1.97), Inches(1.65), Inches(0.36),
     fill=GOLD, line_col=GREEN_DARK, line_w=Pt(1))
txt(s, "★ PATENT APPLIED ★",
    Inches(4.57), Inches(2.0), Inches(1.62), Inches(0.3),
    size=8, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

# ── Slide legend ─────────────────────────────────────────────
legend = [
    (GREEN_DARK, "Field / Hardware"),
    (BLUE,       "Backend / API"),
    (RGBColor(0x7B,0x35,0x00), "Blockchain"),
    (GREEN_MID,  "Consumer Layer"),
]
for i, (col, label) in enumerate(legend):
    lx = Inches(0.28) + i * Inches(3.1)
    rect(s, lx, Inches(4.6), Inches(0.2), Inches(0.2), fill=col)
    txt(s, label, lx+Inches(0.25), Inches(4.59), Inches(2.8), Inches(0.22),
        size=8.5, color=GRAY)


# ── Save ────────────────────────────────────────────────────
out = "/Users/saipranav/Documents/GitHub/SIH-blockchain/HERB_VAULT_CMTI_DIC2026_FULL.pptx"
prs.save(out)
print(f"Saved → {out}")
print(f"Total slides: {len(prs.slides)}")
